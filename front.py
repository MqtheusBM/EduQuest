import streamlit as st
import requests
import json
import PyPDF2
import docx
from pptx import Presentation
import io
import tempfile
import os

def load_css(file_name):
    try:
        with open(file_name) as f:
            st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)
    except FileNotFoundError:
        st.error(f"Ficheiro CSS '{file_name}' não encontrado. Certifica-te de que está na mesma pasta.")

def extrair_texto_arquivo(arquivo):
    texto = ""
    nome_arquivo = arquivo.name.lower()
    
    try:
        arquivo.seek(0)
        
        if nome_arquivo.endswith('.pdf'):
            leitor = PyPDF2.PdfReader(arquivo)
            for pagina in leitor.pages:
                texto_pagina = pagina.extract_text()
                if texto_pagina:
                    texto += texto_pagina + "\n"
                    
        elif nome_arquivo.endswith('.docx') or nome_arquivo.endswith('.pptx'):
            extensao = os.path.splitext(nome_arquivo)[1]
            with tempfile.NamedTemporaryFile(delete=False, suffix=extensao) as tmp:
                tmp.write(arquivo.read())
                tmp_path = tmp.name

            try:
                if nome_arquivo.endswith('.docx'):
                    doc = docx.Document(tmp_path)
                    for paragrafo in doc.paragraphs:
                        texto += paragrafo.text + "\n"
                        
                elif nome_arquivo.endswith('.pptx'):
                    ppt = Presentation(tmp_path)
                    for slide in ppt.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                texto += shape.text + "\n"
            finally:
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
                    
    except Exception as e:
        st.error(f"Erro ao extrair texto do ficheiro: {e}")
        
    return texto

load_css('style.css')


st.title('Bem-vindo ao EduQuest!')
st.write('A Inteligência Artificial que te ajudará a criar as tuas avaliações baseadas nos teus próprios materiais!!')

st.sidebar.title("Configurações da Avaliação")
curso_selecionado = st.sidebar.text_input('Digita o nome do curso:')
dificuldade_selecionada = st.sidebar.selectbox(
    'Seleciona o nível de dificuldade:',
    ('Fácil', 'Médio', 'Difícil'),
)
tipo_questao_selecionado = st.sidebar.selectbox(
    'Seleciona o tipo de questão:',
    ('Múltipla escolha', 'Verdadeiro ou falso', 'Dissertativa'),
)

st.sidebar.markdown('---')
st.sidebar.write('Desenvolvedor: \n\n [Matheus Marques](https://github.com/MqtheusBM)')


tema_avaliacao = st.text_area('Qual é o tema principal? (Opcional se anexares um ficheiro)')
col1, col2 = st.columns([2, 2])
with col1:
    num_questoes = st.number_input('Digita a quantidade de questões:', min_value=1, max_value=50, value=5)
with col2:
    uploaded_file = st.file_uploader("Anexa o teu material base (PDF, PPTX, DOCX)", type=['pdf', 'pptx', 'docx'])


st.markdown("---")
if st.button('Gerar Avaliação'):
    
    contexto_extraido = ""
    if uploaded_file is not None:
        with st.spinner('A extrair texto do ficheiro...'):
            contexto_extraido = extrair_texto_arquivo(uploaded_file)
            st.success("Texto extraído com sucesso!")

    if tema_avaliacao.strip() or contexto_extraido.strip():
        api_url = "http://127.0.0.1:8000/gerar-avaliacao/"

        dados_para_api = {
            "tema": tema_avaliacao.strip(),
            "curso": curso_selecionado,
            "dificuldade": dificuldade_selecionada,
            "tipo_questao": tipo_questao_selecionado,
            "num_questoes": num_questoes,
            "contexto": contexto_extraido
        }

        with st.spinner('Aguarda... a IA está a ler o material e a criar a tua avaliação!'):
            try:
                response = requests.post(api_url, data=json.dumps(dados_para_api))

                if response.status_code == 200:
                    resultado = response.json()
                    st.markdown("### Avaliação Gerada")
                    st.markdown(resultado['avaliacao_gerada'])
                else:
                    st.error(f"Erro ao contactar a API. Código de status: {response.status_code}")
                    st.error(f"Detalhes: {response.text}")

            except requests.exceptions.ConnectionError:
                st.error("Não foi possível conectar à API. Verifica se o servidor FastAPI (uvicorn) está em execução.")

    else:
        st.warning("Por favor, descreve o tema ou anexa um ficheiro de referência antes de gerar a avaliação.")

st.write("O EduQuest pode cometer erros. Por isso, verifica as respostas.")